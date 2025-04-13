import os
import yaml
import logging
from importlib.resources import files
from pptx import Presentation
from pptx.util import Inches, Pt
from .utils import parse_color, get_alignment, get_vertical_anchor

logger = logging.getLogger(__name__)

class PresentationBuilder:
    def __init__(self, template_path, default_slide_layout=6):
        self.prs = Presentation(template_path)
        self.default_slide_layout = default_slide_layout
        self.template_defaults = self._load_layout_defaults()
        self.shape_defaults = self._load_shape_defaults("resources/shape_layouts/shape_defaults.yml")

        self.shape_handlers = {
            "text": self._add_text_shape,
            "table": self._add_table_shape,
            "image": self._add_image_shape
        }

    def _load_layout_defaults(self, layout_path=None):
        if not layout_path:
            return {
                "font": "Avenir LT Next Pro",
                "font_size": 18,
                "left": 0.67,
                "top": 0.4,
                "horizontal_align": "left",
                "vertical_align": "top"
            }
        with open(layout_path, 'r') as f:
            return yaml.safe_load(f).get("defaults", {})

    def _load_shape_defaults(self, resource_path):
        try:
            full_path = files("pptx_layout_engine").joinpath(resource_path)
            with full_path.open("r", encoding="utf-8") as f:
                return yaml.safe_load(f)
        except Exception as e:
            logger.warning(f"Could not load shape defaults from {resource_path}: {e}")
            return {}

    def _resolve_template_layout(self, template_layout):
        if isinstance(template_layout, int):
            return self.prs.slide_layouts[template_layout]
        elif isinstance(template_layout, str):
            for layout in self.prs.slide_layouts:
                if layout.name.strip().lower() == template_layout.strip().lower():
                    return layout
        return self.prs.slide_layouts[self.default_slide_layout]

    def _apply_text_style(self, run, shape_cfg, defaults, color_schemes):
        if "font" in shape_cfg:
            run.font.name = shape_cfg["font"]
        if "font_size" in shape_cfg:
            run.font.size = Pt(shape_cfg["font_size"])
        if shape_cfg.get("font_weight", "").lower() == "bold":
            run.font.bold = True
        if shape_cfg.get("font_italic", False):
            run.font.italic = True
        if "font_color" in shape_cfg:
            run.font.color.rgb = parse_color(shape_cfg["font_color"], color_schemes)

    def _fill_placeholder(self, slide, name, value, shape_cfg, defaults, color_schemes):
        try:
            idx = int(name.replace("placeholder", ""))
            placeholder = slide.placeholders[idx]
        except (ValueError, IndexError):
            logger.warning(f"⚠️ Could not find placeholder '{name}'")
            return

        frame = placeholder.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = get_vertical_anchor(
            shape_cfg.get("vertical_align", defaults.get("vertical_align", "top"))
        )

        if isinstance(value, list):
            for i, line in enumerate(value):
                para = frame.add_paragraph() if i > 0 else frame.paragraphs[0]
                para.text = line
                para.level = 0
                para.alignment = get_alignment(shape_cfg.get("horizontal_align", defaults.get("horizontal_align", "left")))
                self._apply_text_style(para.runs[0], shape_cfg, defaults, color_schemes)
        else:
            para = frame.paragraphs[0]
            para.text = value
            para.alignment = get_alignment(shape_cfg.get("horizontal_align", defaults.get("horizontal_align", "left")))
            self._apply_text_style(para.runs[0], shape_cfg, defaults, color_schemes)

    def _add_text_shape(self, slide, text, shape_cfg, defaults, color_schemes):
        left = Inches(shape_cfg.get("left", defaults["left"]))
        top = Inches(shape_cfg.get("top", defaults["top"]))
        width = Inches(shape_cfg.get("width", 5))
        height = Inches(shape_cfg.get("height", 1))

        textbox = slide.shapes.add_textbox(left, top, width, height)
        frame = textbox.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = get_vertical_anchor(shape_cfg.get("vertical_align", defaults["vertical_align"]))

        if isinstance(text, list):
            for i, line in enumerate(text):
                para = frame.add_paragraph() if i > 0 else frame.paragraphs[0]
                para.text = f"• {line}" if shape_cfg.get("bullet", False) else line
                para.level = 0
                para.alignment = get_alignment(shape_cfg.get("horizontal_align", defaults["horizontal_align"]))
                self._apply_text_style(para.runs[0], shape_cfg, defaults, color_schemes)
        else:
            para = frame.paragraphs[0]
            para.text = text
            para.alignment = get_alignment(shape_cfg.get("horizontal_align", defaults["horizontal_align"]))
            self._apply_text_style(para.runs[0], shape_cfg, defaults, color_schemes)

    def _add_table_shape(self, slide, table_data, shape_cfg, defaults, color_schemes):
        rows, cols = len(table_data), len(table_data[0]) if table_data else 1
        left = Inches(shape_cfg.get("left", defaults["left"]))
        top = Inches(shape_cfg.get("top", defaults["top"]))
        width = Inches(shape_cfg.get("width", 5.5))
        height = Inches(shape_cfg.get("height", 3.0))

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        # Optional fixed row height
        row_height = shape_cfg.get("row_height", None)
        if row_height:
            for i in range(rows):
                table.rows[i].height = Inches(row_height)

        # Optional fixed column widths
        column_widths = shape_cfg.get("column_widths", None)
        if column_widths and len(column_widths) == cols:
            for i, w in enumerate(column_widths):
                table.columns[i].width = Inches(w)

        font_name = shape_cfg.get("font", defaults["font"])
        font_size = Pt(shape_cfg.get("font_size", defaults["font_size"]))
        cell_padding = Inches(shape_cfg.get("cell_padding", 0.05))
        align = get_alignment(shape_cfg.get("horizontal_align", defaults["horizontal_align"]))
        valign = get_vertical_anchor(shape_cfg.get("vertical_align", defaults["vertical_align"]))

        # === Table Color Scheme ===
        header_fill = None
        row_fills = []

        if "table_colors" in shape_cfg:
            scheme_name = shape_cfg["table_colors"]
            scheme = color_schemes.get(scheme_name, [])
            if isinstance(scheme, list):
                for color_dict in scheme:
                    for key, val in color_dict.items():
                        if key.startswith("header"):
                            header_fill = parse_color(val, color_schemes)
                        elif key.startswith("row"):
                            row_fills.append(parse_color(val, color_schemes))

        for r, row_data in enumerate(table_data):
            for c, val in enumerate(row_data):
                cell = table.cell(r, c)
                cell.text = str(val)
                cell.text_frame.word_wrap = True
                cell.text_frame.vertical_anchor = valign
                cell.margin_left = cell.margin_right = cell.margin_top = cell.margin_bottom = cell_padding
                para = cell.text_frame.paragraphs[0]
                para.alignment = align
                run = para.runs[0]
                run.font.name = font_name
                run.font.size = font_size
                if r == 0 and shape_cfg.get("header_bold", True):
                    run.font.bold = True

                if r == 0 and header_fill:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = header_fill
                elif r > 0 and row_fills:
                    fill_color = row_fills[(r - 1) % len(row_fills)]
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = fill_color

    def _add_image_shape(self, slide, image_path, shape_cfg, defaults, color_scheme):
        if not image_path or not os.path.exists(image_path):
            logger.warning(f"Missing image at path: {image_path}")
            return
        left = Inches(shape_cfg.get("left", 1.0))
        top = Inches(shape_cfg.get("top", 1.0))
        width = Inches(shape_cfg.get("width", 4.0))
        height = Inches(shape_cfg.get("height", 3.0))
        slide.shapes.add_picture(image_path, left, top, width, height)

    def clear_all_slides(self):
        while self.prs.slides:
            self.prs.slides._sldIdLst.remove(self.prs.slides._sldIdLst[0])

    def add_slide(self, layout, content):
        with open(layout, 'r') as f:
            layout_cfg = yaml.safe_load(f)

        slide_layout = self._resolve_template_layout(layout_cfg.get("template_layout"))
        slide = self.prs.slides.add_slide(slide_layout)
        layout_defaults = {**self.template_defaults, **layout_cfg.get("defaults", {})}
        color_schemes = layout_cfg.get("color_schemes", {})
        shapes = layout_cfg.get("shapes", {})

        for name, value in content.items():
            layout_shape = shapes.get(name, {})
            default_shape = self.shape_defaults.get(name, {})
            shape_cfg = {**default_shape, **layout_shape}

            if name.startswith("placeholder"):
                self._fill_placeholder(slide, name, value, shape_cfg, layout_defaults, color_schemes)
            else:
                shape_type = shape_cfg.get("type", "text")
                handler = self.shape_handlers.get(shape_type)
                if handler:
                    handler(slide, value, shape_cfg, layout_defaults, color_schemes)

    def build_from_yaml(self, deck_yaml_path, clear_slides=False):
        with open(deck_yaml_path, "r") as f:
            deck = yaml.safe_load(f)

        if template_path := deck.get("template"):
            self.prs = Presentation(template_path)
            if clear_slides:
                self.clear_all_slides()

        defaults = deck.get("defaults", {})
        layout_root = defaults.get("slide_layout_path", "")
        self.template_defaults.update(defaults)

        for idx, slide_def in enumerate(deck.get("slides", [])):
            layout_file = slide_def.get("layout")
            content = slide_def.get("content", {})
            if not layout_file:
                logger.warning(f"[Slide {idx}] Missing layout file. Skipping.")
                continue
            full_layout_path = os.path.join(layout_root, layout_file)
            if not os.path.exists(full_layout_path):
                logger.warning(f"[Slide {idx}] Layout not found: {full_layout_path}. Skipping.")
                continue
            self.add_slide(layout=full_layout_path, content=content)

    # Legacy alias
    load_presentation = build_from_yaml

    def save(self, path):
        self.prs.save(path)
